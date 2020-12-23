from pptx import Presentation as PptxPresentation
from dataclasses import dataclass

@dataclass
class SlideNumber:
    right: float
    bottom: float
    index: int

class Presentation:
    def __init__(self, file_path):
        self._presentation = PptxPresentation(file_path)

    @property
    def slides_titles(self):
        titles = []
        for slide in self._presentation.slides:
            try:
                titles.append(slide.shapes.title.text)
            except:
                titles.append('')
        return titles

    def _find_slide_number(self, slide):
        for shape in slide.shapes:
            if 'Slide Number' in shape.name:
                try:
                    return SlideNumber(
                        index=int(shape.text),
                        right=(shape.left + shape.width) / self._presentation.slide_width,
                        bottom=(shape.top + shape.height) / self._presentation.slide_height
                     )
                except ValueError:
                    return None
        return None

    @property
    def slide_numbers(self):
        numbers = []
        for slide in self._presentation.slides:
            numbers.append(self._find_slide_number(slide))
        return numbers
