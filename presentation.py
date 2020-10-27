from pptx import Presentation as PptxPresentation

class Presentation:
    def __init__(self, file_path):
        self._presentation = PptxPresentation(file_path)

    @property
    def slides_titles(self):
        titles = []
        for slide in self._presentation.slides:
            titles.append(slide.shapes.title.text)
        return titles


# TODO delete (just for test for now)
p = Presentation("test_presentations/good.pptx")
print(p.slides_titles)
